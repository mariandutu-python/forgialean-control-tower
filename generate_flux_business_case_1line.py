# flux_business_case_complete.py - EXCEL + PPTX con RISCHI NC MAGGIORE
import openpyxl
from openpyxl.styles import Font, PatternFill, Border, Side
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_business_case_excel():
    """Excel ORIGINALE + RISCHI NC MAGGIORE"""
    wb = openpyxl.Workbook()
    
    # STILI ORIGINALI
    title_font = Font(bold=True, size=16)
    header_font = Font(bold=True, size=12)
    number_font = Font(bold=True, size=14)
    risk_font = Font(bold=True, size=12, color="FF0000")
    
    # =============================================================================
    # EXECUTIVE - TITOLO + RISCHI ORIGINALI
    # =============================================================================
    ws1 = wb.active
    ws1.title = "EXECUTIVE"
    
    ws1['A1'] = "FLUX SHUTTLE PROFILER - BUSINESS CASE 1 LINEA"
    ws1['A1'].font = title_font
    
    ws1['A3'] = "CERTIFICAZIONI EN9100+ RICHIESTE CLIENTI"
    ws1['A3'].font = header_font
    
    # RISCHI NC MAGGIORE (PAIN POINTS)
    ws1['A5'] = "🚨 RISCHIO NC MAGGIORE: Audit follow-up OBBLIGATORIO"
    ws1['A5'].font = risk_font
    ws1['A6'] = "🚨 BLOCCO PRODUZIONE flussatura critica"
    ws1['A6'].font = risk_font
    ws1['A7'] = "🚨 PERDITA CLIENTI certificati Aerospace/Ferroviario"
    ws1['A7'].font = risk_font
    ws1['A8'] = "Senza sistema = SOSPENSIONE QUALIFICA FORNITORE"
    
    # KPI ORIGINALI
    ws1['A10'] = "INVESTIMENTO TOTALE:"
    ws1['B10'] = 4750
    ws1['A11'] = "RISPARMIO MENSILE:"
    ws1['B11'] = 2170  
    ws1['A12'] = "ROI:"
    ws1['B12'] = "=B10/B11"
    ws1['A13'] = "RISPARMIO ANNUO:"
    ws1['B13'] = "=B11*12"
    
    for row in [10,11,12,13]:
        ws1[f'B{row}'].font = number_font
        ws1[f'B{row}'].number_format = '"€"#,##0'
    ws1['B12'].number_format = '0.0 "mesi"'
    
    ws1['A15'] = "€4.750 = SALVAZIONE da NC MAGGIORE"
    ws1['A15'].font = Font(bold=True, size=14, color="008000")
    
    # =============================================================================
    # COSTI/ROI/PIANO (IDENTICI ORIGINALI)
    # =============================================================================
    ws2 = wb.create_sheet("COSTI")
    ws2['A1'] = "BREAKDOWN INVESTIMENTO €4.750"; ws2['A1'].font = title_font
    
    costi = [
        ["COMPONENTE", "Q.TY", "COSTO €", "NOTE"],
        ["Proxitron 8046AFKM", "3x", 960, "Flow sensor"],
        ["Keyence GT2", "12x", 1200, "Capacitive"],
        ["RPi4 + Stepper", "3x", 345, "Controller"],
        ["Telaio CNC", "3x", 795, "Structure"],
        ["Cavi M12 PTFE", "", 250, ""],
        ["SUBTOTAL HARDWARE", "", "=SUM(C3:C7)", ""],
        ["SOFTWARE + CERTIFICAZIONI", "", "", ""],
        ["Dashboard Streamlit", "30h", 1050, ""],
        ["Database EN9100", "15h", 525, ""],
        ["Test + Validazione", "", 500, ""],
        ["Training", "", 250, ""],
        ["Documentazione", "", 150, ""],
        ["SUBTOTAL SOFTWARE", "", "=SUM(C9:C12)", ""],
        ["GRAND TOTAL", "", "=C8+C13", "CHIUDE NC"]
    ]
    
    for r, row in enumerate(costi, 1):
        for c, value in enumerate(row, 1):
            cell = ws2.cell(row=r+1, column=c); cell.value = value
            if r == 1: cell.font = header_font
            elif c == 3 and isinstance(value, (int, float)): 
                cell.font = number_font; cell.number_format = '"€"#,##0'
    
    # ROI
    ws3 = wb.create_sheet("ROI")
    ws3['A1'] = "RISPARMI MENSILI"; ws3['A1'].font = title_font
    roi = [
        ["VOCE", "PRIMA €", "DOPO €", "RISPARMIO €"],
        ["Validazione", 600, 5, "=B3-C3"],
        ["Controllo manuale", 400, 0, "=B4-C4"],
        ["Difetti", 1000, 200, "=B5-C5"],
        ["Fermo", 200, 50, "=B6-C6"],
        ["Audit", 250, 25, "=B7-C7"],
        ["TOTAL MENSILE", "", "", "=SUM(D3:D7)"],
        ["ROI vs €4.750", 4750, "=D7", "=B8/D8"]
    ]
    
    for r, row in enumerate(roi, 1):
        for c, value in enumerate(row, 1):
            cell = ws3.cell(row=r+1, column=c); cell.value = value
            if r == 1: cell.font = header_font
            elif c in [2,3,4] and r > 1: 
                cell.font = number_font; cell.number_format = '"€"#,##0'
    
    # PIANO
    ws4 = wb.create_sheet("PIANO")
    ws4['A1'] = "6 SETTIMANE - CHIUDE NC MAGGIORE"; ws4['A1'].font = title_font
    piano = [
        ["SETT", "ATTIVITÀ", "COSTO €", "MILESTONE"],
        [1, "Prototipo", 475, "Dashboard live"],
        [2, "Test 5 programmi", 0, "Fluxer OK"],
        [3, "3 telai", 2350, "Linea installata"],
        [4, "SPC + DB", 1050, "Real time 2Hz"],
        [5, "Training + cert.", 400, "EN9100 ready"],
        [6, "Go-Live", 475, "NC MAGGIORE CHIUSO"],
        ["TOT", "", "=SUM(C2:C7)", "CLIENTI SALVATI"]
    ]
    
    for r, row in enumerate(piano, 1):
        for c, value in enumerate(row, 1):
            cell = ws4.cell(row=r+1, column=c); cell.value = value
            if r == 1: cell.font = header_font
            elif c == 3 and isinstance(value, (int, float)): 
                cell.font = number_font; cell.number_format = '"€"#,##0'
    
    # Auto-size sicuro
    for ws in [ws1, ws2, ws3, ws4]:
        for column in ws.columns:
            max_length = max((len(str(cell.value or "")) for cell in column), default=0)
            ws.column_dimensions[column[0].column_letter].width = min(max_length + 2, 25)
    
    filename = f"FluxShuttle_1Line_RISCHI_{datetime.now().strftime('%Y%m%d_%H%M')}.xlsx"
    wb.save(filename)
    return filename

def create_business_case_pptx():
    """PPTX - STILE EXECUTIVE"""
    prs = Presentation()
    
    # SLIDE 1: TITOLO
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide1.shapes.title
    title.text = "FLUX SHUTTLE PROFILER"
    subtitle = slide1.placeholders[1]
    subtitle.text = "BUSINESS CASE 1 LINEA - CHIUDE NC MAGGIORE"
    
    # SLIDE 2: RISCHI (PAIN)
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    title2 = slide2.shapes.title
    title2.text = "🚨 RISCHI NC MAGGIORE EN9100+"
    
    content2 = slide2.placeholders[1]
    content2.text = (
        "❌ AUDIT FOLLOW-UP OBBLIGATORIO (3 mesi)\n"
        "❌ BLOCCO PRODUZIONE flussatura CRITICA\n"
        "❌ PERDITA CLIENTI certificati\n"
        "❌ SOSPENSIONE qualifica fornitore\n"
        "💰 RISCHIO TOTALE: €90K+"
    )
    
    # SLIDE 3: CERTIFICAZIONI
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    title3 = slide3.shapes.title
    title3.text = "CERTIFICAZIONI EN9100+ RICHIESTE CLIENTI"
    content3 = slide3.placeholders[1]
    content3.text = "Senza sistema = SOSPENSIONE QUALIFICA FORNITORE"
    
    # SLIDE 4: KPI
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    title4 = slide4.shapes.title
    title4.text = "€4.750 INVESTIMENTO"
    content4 = slide4.placeholders[1]
    content4.text = (
        "💰 RISPARMIO MENSILE: €2.170\n"
        "🎯 ROI: 2.2 mesi\n"
        "📈 RISPARMIO ANNUO: €26.040\n"
        "⏱️ IMPLEMENTAZIONE: 6 settimane"
    )
    
    # SLIDE 5: PIANO
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    title5 = slide5.shapes.title
    title5.text = "PIANO 6 SETTIMANE"
    content5 = slide5.placeholders[1]
    content5.text = (
        "📅 SETT 1: Prototipo live\n"
        "📅 SETT 3: Linea installata\n"
        "📅 SETT 4: SPC real-time 2Hz\n"
        "📅 SETT 6: NC MAGGIORE CHIUSO"
    )
    
    filename = f"FluxShuttle_1Line_RISCHI_{datetime.now().strftime('%Y%m%d_%H%M')}.pptx"
    prs.save(filename)
    return filename

# ESEGUI ENTRAMBI
if __name__ == "__main__":
    print("🚀 FLUX SHUTTLE BUSINESS CASE + RISCHI NC...")
    excel_file = create_business_case_excel()
    pptx_file = create_business_case_pptx()
    print(f"\n✅ EXCEL: {excel_file}")
    print(f"✅ PPTX:  {pptx_file}")
    print("\n🎯 PRONTO PER DIREZIONE!")
    print("🚨 RISCHI: NC Maggiore | Blocco produzione | Perdita clienti")