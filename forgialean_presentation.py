from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import matplotlib.pyplot as plt
import pandas as pd
from io import BytesIO
import os

# COLORI BRAND FORGIALEAN
COLOR_PRIMARY = RGBColor(0, 102, 204)  # Blu professionale
COLOR_SECONDARY = RGBColor(220, 53, 69)  # Rosso per alert
COLOR_SUCCESS = RGBColor(40, 167, 69)  # Verde per target
COLOR_LIGHT = RGBColor(240, 240, 240)  # Grigio chiaro
COLOR_DARK = RGBColor(33, 37, 41)  # Grigio scuro

# Percorso logo
LOGO_PATH = r"C:\Users\marian.dutu\Desktop\ForgiaLean\forgialean_logo.png"

def add_title_slide_pro(prs, title_text, subtitle_text=""):
    """Slide titolo professionale con logo e background"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_PRIMARY
    
    # Logo (se esiste)
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, Inches(4.2), Inches(0.5), width=Inches(1.5))
    
    # Titolo
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title_text
    title_frame.paragraphs[0].font.size = Pt(54)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Sottotitolo
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.9), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle_text
    subtitle_frame.paragraphs[0].font.size = Pt(26)
    subtitle_frame.paragraphs[0].font.color.rgb = COLOR_LIGHT
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Footer
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.5))
    footer_frame = footer_box.text_frame
    footer_frame.text = "ForgiaLean – Supporto alle PMI manifatturiere con OEE < 80%"
    footer_frame.paragraphs[0].font.size = Pt(13)
    footer_frame.paragraphs[0].font.color.rgb = COLOR_LIGHT
    footer_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title_text, content_type="bullets", data=None):
    """Slide contenuto con design professionale"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Background bianco
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Header bar blu
    header_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.9))
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = COLOR_PRIMARY
    header_shape.line.color.rgb = COLOR_PRIMARY
    
    # Logo piccolo in alto a sinistra
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, Inches(0.3), Inches(0.15), width=Inches(0.6))
    
    # Titolo
    title_box = slide.shapes.add_textbox(Inches(1.2), Inches(0.15), Inches(8.3), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = title_text
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # Contenuto
    if content_type == "bullets":
        body_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8.4), Inches(5.8))
        text_frame = body_box.text_frame
        text_frame.word_wrap = True
        for i, bullet in enumerate(data):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            p.text = bullet
            p.level = 0
            p.font.size = Pt(18)
            p.font.color.rgb = COLOR_DARK
            p.space_before = Pt(8)
    
    elif content_type == "image":
        slide.shapes.add_picture(data, Inches(1), Inches(1.2), width=Inches(8))
    
    # Footer
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(7.2), Inches(9), Inches(0.2))
    footer_frame = footer_box.text_frame
    footer_frame.text = f"ForgiaLean | Slide {len(prs.slides)}"
    footer_frame.paragraphs[0].font.size = Pt(10)
    footer_frame.paragraphs[0].font.color.rgb = COLOR_LIGHT
    
    return slide

def create_oee_trend_chart():
    """Crea grafico trend OEE"""
    mesi = ['Gen', 'Feb', 'Mar', 'Apr', 'Mag', 'Giu']
    oee_values = [72, 74, 76, 78, 81, 84]
    
    fig, ax = plt.subplots(figsize=(10, 6))
    
    ax.plot(mesi, oee_values, marker='o', linewidth=3, markersize=10, 
            color='#0066CC', label='OEE Attuale')
    ax.axhline(y=80, color='#DC3545', linestyle='--', linewidth=2, label='Target 80%')
    ax.fill_between(range(len(mesi)), 0, oee_values, alpha=0.2, color='#0066CC')
    ax.fill_between(range(len(mesi)), 80, 100, alpha=0.1, color='#28A745')
    
    ax.set_ylabel('OEE (%)', fontsize=12, fontweight='bold')
    ax.set_xlabel('Periodo', fontsize=12, fontweight='bold')
    ax.set_ylim(60, 100)
    ax.grid(True, alpha=0.3)
    ax.legend(fontsize=11, loc='lower right')
    ax.set_title('Trend OEE - Ultimi 6 mesi', fontsize=14, fontweight='bold')
    
    img_stream = BytesIO()
    plt.savefig(img_stream, format='png', dpi=300, bbox_inches='tight')
    img_stream.seek(0)
    plt.close()
    
    return img_stream

def create_pain_pareto_chart():
    """Crea Pareto pain"""
    causes = ['Fermi impianto', 'Scarti qualità', 'Rallentamenti', 'Setup']
    values = [45, 25, 20, 10]
    
    fig, ax = plt.subplots(figsize=(10, 6))
    
    bars = ax.bar(causes, values, color=['#DC3545', '#FF6B35', '#FFC300', '#28A745'], 
                   edgecolor='black', linewidth=1.5)
    
    total = sum(values)
    for bar, val in zip(bars, values):
        height = bar.get_height()
        ax.text(bar.get_x() + bar.get_width()/2., height,
                f'{val}%', ha='center', va='bottom', fontsize=11, fontweight='bold')
    
    ax.set_ylabel('Impatto sulla perdita OEE (%)', fontsize=12, fontweight='bold')
    ax.set_title('Analisi Pareto - Cause principali di perdita OEE', fontsize=14, fontweight='bold')
    ax.set_ylim(0, 55)
    ax.grid(True, alpha=0.3, axis='y')
    
    img_stream = BytesIO()
    plt.savefig(img_stream, format='png', dpi=300, bbox_inches='tight')
    img_stream.seek(0)
    plt.close()
    
    return img_stream

def create_kpi_breakdown_chart():
    """Crea breakdown OEE"""
    categories = ['Disponibilità', 'Performance', 'Qualità']
    values = [90, 87, 80]
    
    fig, ax = plt.subplots(figsize=(10, 6))
    
    colors = ['#28A745' if v >= 80 else '#DC3545' for v in values]
    bars = ax.bar(categories, values, color=colors, edgecolor='black', linewidth=2, width=0.6)
    
    for bar, val in zip(bars, values):
        height = bar.get_height()
        ax.text(bar.get_x() + bar.get_width()/2., height,
                f'{val}%', ha='center', va='bottom', fontsize=16, fontweight='bold')
    
    ax.set_ylabel('Componente OEE (%)', fontsize=12, fontweight='bold')
    ax.set_title('Breakdown OEE - Componenti principali', fontsize=14, fontweight='bold')
    ax.set_ylim(0, 100)
    ax.axhline(y=80, color='#DC3545', linestyle='--', linewidth=2, alpha=0.5)
    ax.grid(True, alpha=0.3, axis='y')
    
    img_stream = BytesIO()
    plt.savefig(img_stream, format='png', dpi=300, bbox_inches='tight')
    img_stream.seek(0)
    plt.close()
    
    return img_stream

def crea_presentazione_forgialean_pro(output_filename="ForgiaLean_Presentazione_Pro.pptx"):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # SLIDE 1 – TITOLO
    add_title_slide_pro(
        prs,
        "ForgiaLean",
        "Quando l'OEE fa male – Supporto alle PMI manifatturiere"
    )
    
    # SLIDE 2 – APERTURA
    add_content_slide(
        prs,
        "Conosci l'OEE della tua linea principale?",
        "bullets",
        [
            "📊 Superiore all'85%? → Buon posizionamento, mantenere.",
            "📊 Tra 80% e 85%? → Margini di miglioramento visibili.",
            "📊 Tra 70% e 80%? → Potenziale di recupero significativo.",
            "📊 Sotto il 70%? → Situazione critica, intervento urgente.",
            "",
            "💡 Se non conosci il valore, questo è il primo campanello d'allarme."
        ]
    )
    
    # SLIDE 3 – SEGNALI D'ALLARME
    add_content_slide(
        prs,
        "Segnali che l'OEE sta facendo male",
        "bullets",
        [
            "🔴 Fermi impianto ricorrenti senza visione chiara delle cause.",
            "🔴 Scarti di qualità gestiti solo a consuntivo, non prevenuti.",
            "🔴 Performance macchina lontana dai valori di targa.",
            "🔴 Report in Excel che richiedono ore, ma danno poche risposte.",
            "🔴 Impossibilità di identificare le 3 principali fonti di perdita."
        ]
    )
    
    # SLIDE 4 – TREND OEE
    slide = add_content_slide(prs, "Trend OEE - Esempio di intervento", "image", None)
    oee_chart = create_oee_trend_chart()
    slide.shapes.add_picture(oee_chart, Inches(0.5), Inches(1.2), width=Inches(9))
    
    # SLIDE 5 – PARETO
    slide = add_content_slide(prs, "Analisi Pareto - Dove agire", "image", None)
    pareto_chart = create_pain_pareto_chart()
    slide.shapes.add_picture(pareto_chart, Inches(0.5), Inches(1.2), width=Inches(9))
    
    # SLIDE 6 – BREAKDOWN
    slide = add_content_slide(prs, "Breakdown OEE - Componenti", "image", None)
    breakdown_chart = create_kpi_breakdown_chart()
    slide.shapes.add_picture(breakdown_chart, Inches(0.5), Inches(1.2), width=Inches(9))
    
    # SLIDE 7 – OFFERTA
    add_content_slide(
        prs,
        "Cosa offre ForgiaLean",
        "bullets",
        [
            "✅ Analisi approfondita dei dati di produzione.",
            "✅ Identificazione delle 3 principali fonti di perdita OEE.",
            "✅ Dashboard operative per il monitoraggio quotidiano.",
            "✅ Definizione di azioni concrete su disponibilità, performance, qualità.",
            "✅ Supporto strutturato verso target OEE > 80%."
        ]
    )
    
    # SLIDE 8 – APPROCCIO
    add_content_slide(
        prs,
        "Approccio ForgiaLean - Interattivo",
        "bullets",
